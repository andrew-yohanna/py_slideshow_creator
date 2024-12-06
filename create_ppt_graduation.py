import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
from collections import Counter
from pptx.enum.text import PP_ALIGN

# Folder containing all graduate folders
root_folder = os.getcwd() + "/All"
presentation_path = "Graduates_Presentation.pptx"  # Path to the output presentation
# Check if the presentation already exists and delete it
if os.path.exists(presentation_path):
    os.remove(presentation_path)
    print(f"Deleted existing presentation: {presentation_path}")

presentation = Presentation()

def get_dominant_color(image_path):
    """
    Extracts the most dominant color from an image.
    """
    with Image.open(image_path) as img:
        img = img.resize((100, 100))  # Resize for faster processing
        pixels = img.getdata()
        pixels = [pixel[:3] for pixel in pixels]  # Exclude alpha channel if present
        most_common = Counter(pixels).most_common(1)[0][0]  # Get the most common color
        return most_common
    
def get_contrast_color(rgb):
    """
    Returns a high-contrast color (either black or white) for the given RGB color.
    """
    r, g, b = rgb
    luminance = 0.2126 * r + 0.7152 * g + 0.0722 * b
    return RGBColor(255, 255, 255) if luminance < 128 else RGBColor(0, 0, 0)

def find_photo(graduate_folder, base_filename):
    """
    Returns the path of the first photo found with the given base filename and possible extensions.
    """
    possible_extensions = ['.jpg', '.jpeg', '.png']  # List of extensions to check

    for ext in possible_extensions:
        photo_path = os.path.join(graduate_folder, base_filename + ext)
        if os.path.exists(photo_path):
            return photo_path
    return None  # Return None if no valid photo is found

def add_slide_with_photo(presentation, graduate_name, photo_path, label):
    """
    Adds a slide with a photo filling the slide and text at the top with a background color.
    """
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # Blank slide layout

    # Get the dominant color
    dominant_color = get_dominant_color(photo_path)
    rgb_color = RGBColor(*dominant_color)
    contrast_color = get_contrast_color(dominant_color)

    # Add the photo
    if os.path.exists(photo_path):
        pic = slide.shapes.add_picture(photo_path, Inches(0), Inches(1), width=presentation.slide_width)
        # Center the image vertically and horizontally
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height
        pic.left = (slide_width - pic.width) // 2
        pic.top = (slide_height - pic.height) // 2
 
    slide_label = f"{graduate_name} - {label}"

    # Add name and label text at the bottom-right corner
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height

    # Define the size and position of the text box
    textbox_width = Inches(2) + Inches(0.1 * len(slide_label))  # Width of the text box
    textbox_height = Inches(1)  # Height of the text box

    # Place the text box in the bottom-right corner
    left = slide_width - textbox_width - Inches(0.5)  # Right alignment with margin
    top = slide_height - textbox_height - Inches(0.5)  # Bottom alignment with margin
    
    # Add name and label text at the top
    # title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(1))
    title_box = slide.shapes.add_textbox(left, top, textbox_width, textbox_height)
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = slide_label
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = contrast_color  # Set the text color for contrast
   
    # Center the text horizontally
    p.alignment = PP_ALIGN.CENTER

    # Add border to the text box
    line = title_box.line  # Get the line (border) object
    line.color.rgb = contrast_color  # Set the border color to black
    line.width = Inches(0.05)  # Set the width of the border (can adjust)

    # Set background color for the text box
    fill = title_box.fill
    fill.solid()
    fill.fore_color.rgb = rgb_color

for folder in os.listdir(root_folder):
    graduate_folder = os.path.join(root_folder, folder)
    if os.path.isdir(graduate_folder):
        graduate_name = folder  # Graduate's name (folder name)
        # recent_photo = os.path.join(graduate_folder, "recent.jpg")
        # younger_photo = os.path.join(graduate_folder, "baby.jpg")

        # Get the recent and younger photos (this works for any supported extension)
        recent_photo = find_photo(graduate_folder, "recent")
        younger_photo = find_photo(graduate_folder, "baby")

        if younger_photo:
            add_slide_with_photo(presentation, graduate_name, younger_photo, "Then")
        else:
            print(f"Baby photo missing for: {graduate_name}")

        # Print the results (for debugging)
        if recent_photo:
            add_slide_with_photo(presentation, graduate_name, recent_photo, "Now")
        else:
             print(f"Recent photo missing for: {graduate_name}")

        # # Add slides for each photo
        # if os.path.exists(recent_photo):
        #     add_slide_with_photo(presentation, graduate_name, recent_photo, "Now")
        # else:
        #     print(f"Now photo missing for: {graduate_name}")

        # if os.path.exists(younger_photo):
        #     add_slide_with_photo(presentation, graduate_name, younger_photo, "Younger")
        # else:
        #     print(f"Younger photo missing for: {graduate_name}")

# Save the new presentation
presentation.save(presentation_path)
print(f"New presentation saved as: {presentation_path}")
